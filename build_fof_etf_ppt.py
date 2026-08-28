#!/usr/bin/env python3
"""Build FOF-ETF PPT without deepcopy (PowerPoint-compatible)."""
from __future__ import annotations

import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

BASE = "/home/ubuntu/.cursor/projects/workspace/uploads/FOF-ETF_______fe89.pptx"
ASSETS = "/workspace/ppt_assets"
OUTPUT = "/workspace/FOF-ETF-ppt/FOF-ETF产品_吴致超_完整版.pptx"

DISCLAIMER = "基金有风险，投资需谨慎。禁止第三方机构摘引、截取或以其他不恰当方式转播。"

# Standard positions from original deck (EMU)
TITLE_BOX = (513430, 318990, 7795192, 430887)
CHART_BOX = (614098, 1415716, 9663764, 4026568)
NOTE_BOX = (336332, 6585306, 8841122, 252730)
DISC_BOX = (226244, 6556585, 6096000, 245110)
LOGO_BOX = (10557091, 6240269, 1136741, 290354)


def reorder_slides(pres: Presentation, order: list[int]):
    sldIdLst = pres.slides._sldIdLst
    slides = list(sldIdLst)
    ordered = [slides[i] for i in order]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for child in ordered:
        sldIdLst.append(child)


def set_title_on_slide(slide, title: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if not t or DISCLAIMER in t or t.startswith("注：") or t.startswith("数据来源"):
            continue
        if len(t) < 90 and ("：" in t or t in ("投资体系", "团队介绍", "资产配置", "ETF轮动", "ETF收益增强")):
            shape.text_frame.text = title
            return


def replace_largest_picture(slide, image_path: str):
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    content_pics = sorted(
        [p for p in pics if p.width > Inches(2.5)],
        key=lambda p: p.width * p.height,
        reverse=True,
    )
    if not content_pics:
        return
    pic = content_pics[0]
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    pic._element.getparent().remove(pic._element)
    slide.shapes.add_picture(image_path, left, top, width, height)


def _style_title_frame(tf, text: str):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)


def _add_disclaimer(slide):
    box = slide.shapes.add_textbox(*DISC_BOX)
    p = box.text_frame.paragraphs[0]
    p.text = DISCLAIMER
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(128, 128, 128)


def _add_note(slide, note: str):
    box = slide.shapes.add_textbox(*NOTE_BOX)
    p = box.text_frame.paragraphs[0]
    p.text = note
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(89, 89, 89)


def add_content_image_slide(pres: Presentation, title: str, image: str, note: str = ""):
    """Create a new slide using blank layout + native add_picture (no deepcopy)."""
    slide = pres.slides.add_slide(pres.slide_layouts[11])  # 1_标题幻灯片 (empty)
    title_box = slide.shapes.add_textbox(*TITLE_BOX)
    _style_title_frame(title_box.text_frame, title)
    l, t, w, h = CHART_BOX
    slide.shapes.add_picture(image, l, t, w, h)
    if note:
        _add_note(slide, note)
    _add_disclaimer(slide)
    return slide


def add_content_text_slide(pres: Presentation, title: str, bullets: list[str]):
    slide = pres.slides.add_slide(pres.slide_layouts[11])
    title_box = slide.shapes.add_textbox(*TITLE_BOX)
    _style_title_frame(title_box.text_frame, title)
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12.0), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(31, 78, 121)
        p.space_after = Pt(10)
    _add_disclaimer(slide)
    return slide


def add_section_slide(pres: Presentation, number: str, title: str):
    slide = pres.slides.add_slide(pres.slide_layouts[2])  # 节标题
    if slide.shapes.title:
        slide.shapes.title.text = title
    for shape in slide.placeholders:
        if shape != slide.shapes.title and shape.has_text_frame:
            shape.text_frame.text = number
            break
    _add_disclaimer(slide)
    return slide


def patch_existing(prs: Presentation):
    set_title_on_slide(prs.slides[8], "全流程定量定性相结合的投资策略池：资产配置")

    s14 = prs.slides[13]
    for shape in s14.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if "得出基金产品之后" in txt:
            shape.text_frame.text = (
                "基金组合行业配置线上限25%；ETF组合行业配置上限35%。"
                "ETF缺乏主动管理alpha，需适度行业超配以获取超额，同时控制偏离。"
            )
        elif "使用量化模型选出前20只" in txt:
            shape.text_frame.text = "1. 扩展窗口机器学习选基模型，每期筛选Top20基金"
        elif "使用基金行业高频仓位测算的结果" in txt:
            shape.text_frame.text = "2. 穿透测算组合行业权重，映射为ETF轮动配置"

    replace_largest_picture(prs.slides[12], f"{ASSETS}/930950_sector_allocation.png")
    replace_largest_picture(prs.slides[14], f"{ASSETS}/fund_model_nav.png")

    toc = prs.slides[1]
    for shape in toc.shapes:
        if shape.has_text_frame and "投资体系" in shape.text_frame.text:
            shape.text_frame.text = shape.text_frame.text.replace(
                "投资体系", "投资体系（资产配置 / ETF轮动 / ETF收益增强）"
            )


def main():
    shutil.copy(BASE, OUTPUT)
    prs = Presentation(OUTPUT)
    patch_existing(prs)

    base = len(prs.slides)

    add_section_slide(prs, "02", "资产配置")
    add_section_slide(prs, "03", "ETF轮动")
    add_content_text_slide(
        prs,
        "ETF轮动策略：核心思路",
        [
            "ETF轮动策略的本质：为ETF产品配置相应权重，以战胜既定基准。",
            "核心难点：每一期如何决定ETF的配置比例。",
            "银华FOF方案：让最优秀的基金组合指引ETF配置比例，解决最难的环节。",
            "在选基模型已实盘验证稳定超额（相对930950）的基础上，",
            "开发ETF替代组合，回测对930950同样具备稳定超额，具备产品化条件。",
        ],
    )
    add_content_image_slide(
        prs,
        "机器学习选基模型：评价指标",
        f"{ASSETS}/fund_model_kpi.png",
        "注：数据来源银华基金、Wind；行业配置上限25%。",
    )
    add_content_image_slide(
        prs,
        "机器学习选基模型：逐年超额收益",
        f"{ASSETS}/fund_model_annual_excess.png",
        "注：2017-2026年，相对930950多数年份取得正超额。",
    )
    add_content_image_slide(
        prs,
        "ETF替代组合结构：85/10/5",
        f"{ASSETS}/portfolio_structure.png",
        "注：核心85%+卫星10%+货币5%；仅持有行业ETF与货币ETF。",
    )
    add_content_image_slide(
        prs,
        "ETF轮动策略：实施流程",
        f"{ASSETS}/etf_strategy_flow.png",
        "注：扩展窗口选基→行业穿透与35%上限优化→1B规则选ETF→组合落地。",
    )
    add_content_image_slide(
        prs,
        "ETF替代策略：回测核心指标",
        f"{ASSETS}/etf_kpi_table.png",
        "注：扩展窗口|上限35%|85/10/5|1B穿透；单只ETF≤20%。",
    )
    add_content_image_slide(
        prs,
        "ETF替代策略：逐年超额收益（2021起）",
        f"{ASSETS}/etf_annual_excess_compare.png",
        "注：固定窗口与扩展窗口均对930950具备正超额。",
    )
    add_content_text_slide(
        prs,
        "ETF替代策略：操作要点",
        [
            "选基模型：扩展窗口-30天间隔，每期Top20基金，单行业≤35%优化。",
            "核心袖套（85%）：优化后行业权重 → 1B规则映射行业ETF（250日收益优选）。",
            "卫星袖套（10%）：TOP20穿透行业权重 → 同一套ETF映射规则。",
            "货币袖套（5%）：511990.SH，满足开放式基金现金约束。",
            "FOF合规：单只ETF持仓上限20%，超额权重溢向次优候选。",
            "调仓频率约30天；组合层仅行业ETF+货币ETF，不持有主动基金。",
        ],
    )
    add_content_text_slide(
        prs,
        "策略产品化与实盘可行性",
        [
            "选基模型：已实盘检验，相对930950年化超额约10%，超额回撤修复能力强。",
            "ETF替代组合：回测CAGR超额+2.07%（全样本），2021起超额+4.92%。",
            "持仓透明：全部为ETF，便于券商FOF产品发行与日常运作。",
            "规则清晰：调仓日、行业权重、ETF映射均有完整数据链路，可系统化执行。",
            "风控完备：行业上限35%、单券上限20%、现金比例5%，符合FOF合规要求。",
        ],
    )
    add_section_slide(prs, "04", "ETF收益增强")

    n = base
    order = [
        0, 1, 2, 3, 4, 5, 6, 7,
        n, 8, 9, 10,
        n + 1, n + 2, 11, 14, n + 3, n + 4, 15, 13, 12,
        n + 5, n + 6, n + 7, n + 8, n + 9, n + 10,
        16, n + 11, 17, 18, 19, 20, 21,
    ]
    reorder_slides(prs, order)
    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
